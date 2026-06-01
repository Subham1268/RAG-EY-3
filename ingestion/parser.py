"""
ingestion/parser.py
────────────────────
Multimodal document parser supporting PDF, PPTX, DOCX, and XLSX.

For each document it produces:
  - A list of TextChunk objects (page / slide / sheet level)
  - A list of ImageChunk objects (embedded diagrams, charts)
  - A list of TableChunk objects (structured tabular data)

Design notes:
  • PDF  → PyMuPDF for text+images, pdfplumber for precise tables
  • PPTX → python-pptx; each slide = 1 text chunk + extracted shape images
  • DOCX → python-docx; paragraphs + embedded tables
  • XLSX → openpyxl; each sheet → markdown table + numeric summary
"""

from __future__ import annotations

import base64
import io
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

import fitz  # PyMuPDF
import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

# ── Data models ────────────────────────────────────────────────────────────────

DocType = Literal["pdf", "pptx", "docx", "xlsx"]


@dataclass
class TextChunk:
    text: str
    doc_type: DocType
    source_file: str
    page_or_slide: int          # 0-indexed
    section_title: str = ""
    metadata: dict = field(default_factory=dict)


@dataclass
class ImageChunk:
    image_b64: str              # Base64-encoded PNG
    caption: str                # Extracted caption / surrounding text
    doc_type: DocType
    source_file: str
    page_or_slide: int
    metadata: dict = field(default_factory=dict)


@dataclass
class TableChunk:
    markdown: str               # Table rendered as markdown
    doc_type: DocType
    source_file: str
    page_or_slide: int
    metadata: dict = field(default_factory=dict)


@dataclass
class ParsedDocument:
    source_file: str
    doc_type: DocType
    text_chunks: list[TextChunk] = field(default_factory=list)
    image_chunks: list[ImageChunk] = field(default_factory=list)
    table_chunks: list[TableChunk] = field(default_factory=list)


# ── Main parser ────────────────────────────────────────────────────────────────

class DocumentParser:
    """
    Unified entry point. Call `parse(path)` for any supported file type.
    """

    MIN_IMAGE_SIZE_PX = 80          # Skip tiny icons / bullets
    MIN_TEXT_LENGTH   = 30          # Skip near-empty text blocks

    def parse(self, path: str | Path) -> ParsedDocument:
        path = Path(path)
        ext  = path.suffix.lower().lstrip(".")

        dispatch = {
            "pdf":  self._parse_pdf,
            "pptx": self._parse_pptx,
            "docx": self._parse_docx,
            "xlsx": self._parse_xlsx,
        }
        if ext not in dispatch:
            raise ValueError(f"Unsupported file type: {ext}")

        return dispatch[ext](path)

    # ── PDF ────────────────────────────────────────────────────────────────────

    def _parse_pdf(self, path: Path) -> ParsedDocument:
        doc = ParsedDocument(source_file=str(path), doc_type="pdf")

        # ─ Text + embedded images via PyMuPDF ──────────────────────────────
        fitz_doc = fitz.open(str(path))
        for page_idx, page in enumerate(fitz_doc):
            text = page.get_text("text").strip()
            if len(text) >= self.MIN_TEXT_LENGTH:
                doc.text_chunks.append(TextChunk(
                    text=text,
                    doc_type="pdf",
                    source_file=str(path),
                    page_or_slide=page_idx,
                ))

            # Extract embedded raster images
            for img_meta in page.get_images(full=True):
                # xref   = img_meta[0]
                # base_img = fitz_doc.extract_image(xref)
                # w, h     = base_img["width"], base_img["height"]
                # if w < self.MIN_IMAGE_SIZE_PX or h < self.MIN_IMAGE_SIZE_PX:
                #     continue
                # img_b64 = base64.b64encode(base_img["image"]).decode()
                # doc.image_chunks.append(ImageChunk(
                #     image_b64=img_b64,
                #     caption=f"Image on page {page_idx + 1}",
                #     doc_type="pdf",
                #     source_file=str(path),
                #     page_or_slide=page_idx,
                # ))
                pass
        fitz_doc.close()

        # ─ Tables via pdfplumber ───────────────────────────────────────────
        with pdfplumber.open(str(path)) as pdf:
            for page_idx, page in enumerate(pdf.pages):
                for table in page.extract_tables():
                    md = self._table_to_markdown(table)
                    if md:
                        doc.table_chunks.append(TableChunk(
                            markdown=md,
                            doc_type="pdf",
                            source_file=str(path),
                            page_or_slide=page_idx,
                        ))

        return doc

    # ── PPTX ───────────────────────────────────────────────────────────────────

    def _parse_pptx(self, path: Path) -> ParsedDocument:
        doc  = ParsedDocument(source_file=str(path), doc_type="pptx")
        prs  = Presentation(str(path))

        for slide_idx, slide in enumerate(prs.slides):
            slide_texts: list[str] = []
            slide_title = ""

            for shape in slide.shapes:
                # ─ Text frames ─────────────────────────────────────────────
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        if not slide_title and shape.shape_type == 13:   # title
                            slide_title = text
                        slide_texts.append(text)

                # ─ Tables ──────────────────────────────────────────────────
                if shape.has_table:
                    rows = [
                        [cell.text for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    md = self._table_to_markdown(rows)
                    if md:
                        doc.table_chunks.append(TableChunk(
                            markdown=md,
                            doc_type="pptx",
                            source_file=str(path),
                            page_or_slide=slide_idx,
                        ))

                # ─ Images / diagrams ───────────────────────────────────────
                if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_bytes = shape.image.blob
                        img       = Image.open(io.BytesIO(img_bytes))
                        if img.width < self.MIN_IMAGE_SIZE_PX or img.height < self.MIN_IMAGE_SIZE_PX:
                            continue
                        buf = io.BytesIO()
                        img.save(buf, format="PNG")
                        img_b64 = base64.b64encode(buf.getvalue()).decode()
                        doc.image_chunks.append(ImageChunk(
                            image_b64=img_b64,
                            caption=f"Slide {slide_idx + 1}: {slide_title}",
                            doc_type="pptx",
                            source_file=str(path),
                            page_or_slide=slide_idx,
                        ))
                    except Exception:
                        pass

            combined_text = "\n".join(slide_texts)
            if len(combined_text) >= self.MIN_TEXT_LENGTH:
                doc.text_chunks.append(TextChunk(
                    text=combined_text,
                    doc_type="pptx",
                    source_file=str(path),
                    page_or_slide=slide_idx,
                    section_title=slide_title,
                ))

        return doc

    # ── DOCX ───────────────────────────────────────────────────────────────────

    def _parse_docx(self, path: Path) -> ParsedDocument:
        doc     = ParsedDocument(source_file=str(path), doc_type="docx")
        docx    = DocxDocument(str(path))
        current_section: list[str] = []
        current_heading = ""

        for block in docx.element.body:
            tag = block.tag.split("}")[-1]

            if tag == "p":
                from docx.oxml.ns import qn
                para_elem = block
                # Detect heading style
                pPr   = para_elem.find(qn("w:pPr"))
                style = ""
                if pPr is not None:
                    pStyle = pPr.find(qn("w:pStyle"))
                    if pStyle is not None:
                        style = pStyle.get(qn("w:val"), "")

                text = "".join(r.text or "" for r in para_elem.findall(
                    ".//" + qn("w:t")
                )).strip()

                if style.startswith("Heading") and text:
                    # Flush accumulated text
                    self._flush_docx_section(doc, current_section, current_heading, path)
                    current_section = []
                    current_heading = text
                elif text:
                    current_section.append(text)

            elif tag == "tbl":
                # Extract table rows
                from docx.oxml.ns import qn
                rows = []
                for tr in block.findall(".//" + qn("w:tr")):
                    cells = [
                        "".join(t.text or "" for t in tc.findall(".//" + qn("w:t")))
                        for tc in tr.findall(".//" + qn("w:tc"))
                    ]
                    rows.append(cells)
                md = self._table_to_markdown(rows)
                if md:
                    doc.table_chunks.append(TableChunk(
                        markdown=md,
                        doc_type="docx",
                        source_file=str(path),
                        page_or_slide=0,
                    ))

        self._flush_docx_section(doc, current_section, current_heading, path)
        return doc

    def _flush_docx_section(
        self,
        doc: ParsedDocument,
        lines: list[str],
        heading: str,
        path: Path,
    ) -> None:
        text = "\n".join(lines).strip()
        if len(text) >= self.MIN_TEXT_LENGTH:
            doc.text_chunks.append(TextChunk(
                text=text,
                doc_type="docx",
                source_file=str(path),
                page_or_slide=0,
                section_title=heading,
            ))

    # ── XLSX ───────────────────────────────────────────────────────────────────

    def _parse_xlsx(self, path: Path) -> ParsedDocument:
        doc = ParsedDocument(source_file=str(path), doc_type="xlsx")
        wb  = load_workbook(str(path), read_only=True, data_only=True)

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws   = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            # Convert to string rows (handle None)
            str_rows = [
                [str(cell) if cell is not None else "" for cell in row]
                for row in rows
            ]

            # Text chunk: sheet name + first 5 rows as context
            preview = "\n".join(
                " | ".join(r) for r in str_rows[:5]
            )
            doc.text_chunks.append(TextChunk(
                text=f"Sheet: {sheet_name}\n\n{preview}",
                doc_type="xlsx",
                source_file=str(path),
                page_or_slide=sheet_idx,
                section_title=sheet_name,
            ))

            # Full table as markdown (cap at 200 rows for practicality)
            md = self._table_to_markdown(str_rows[:200])
            if md:
                doc.table_chunks.append(TableChunk(
                    markdown=md,
                    doc_type="xlsx",
                    source_file=str(path),
                    page_or_slide=sheet_idx,
                    metadata={"sheet_name": sheet_name},
                ))

        wb.close()
        return doc

    # ── Helpers ────────────────────────────────────────────────────────────────

    @staticmethod
    def _table_to_markdown(rows: list[list]) -> str:
        if not rows:
            return ""
        rows = [[str(c) for c in row] for row in rows]
        header    = rows[0]
        separator = ["---"] * len(header)
        body      = rows[1:]
        lines     = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ] + [
            "| " + " | ".join(row) + " |"
            for row in body
        ]
        return "\n".join(lines)
